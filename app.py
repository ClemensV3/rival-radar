import streamlit as st
import google.generativeai as genai
import pandas as pd
import json
import io
import os
import re
import time
import altair as alt
import feedparser
import urllib.parse
import urllib.request
from youtube_transcript_api import YouTubeTranscriptApi
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

st.set_page_config(page_title="RivalRadar", layout="wide")

st.markdown("""
<style>
    /* Global Font Size Boost */
    html, body, [class*="st-"] {
        font-size: 1.1rem;
    }

    /* Custom Radar Animation - SCALED UP */
    @keyframes sweep {
        0% { transform: rotate(0deg); }
        100% { transform: rotate(360deg); }
    }
    .header-link {
        text-decoration: none;
        display: block;
        cursor: pointer;
    }
    .radar-container {
        display: flex;
        justify-content: center;
        align-items: center;
        margin-bottom: 5px;
        transition: transform 0.2s ease-in-out;
    }
    .radar-container:hover {
        transform: scale(1.02);
    }
    .radar {
        position: relative;
        width: 60px; /* RIESIG */
        height: 60px; /* RIESIG */
        border-radius: 50%;
        border: 3px solid #E3000F; /* SANY Red */
        background: radial-gradient(circle, rgba(227,0,15,0.1) 0%, rgba(26,28,30,0) 70%);
        margin-right: 20px;
        overflow: hidden;
    }
    .radar::before {
        content: '';
        position: absolute;
        top: 50%;
        left: 50%;
        width: 30px;
        height: 3px;
        background-color: #E3000F;
        transform-origin: 0% 50%;
        animation: sweep 2s linear infinite;
    }
    .radar::after {
        content: '';
        position: absolute;
        top: 50%;
        left: 50%;
        transform: translate(-50%, -50%);
        width: 8px;
        height: 8px;
        background-color: #ffffff;
        border-radius: 50%;
    }
    
    /* Typography Overrides */
    .title-text {
        font-size: 4.5rem; /* RIESIG */
        font-weight: 900;
        letter-spacing: 3px;
        color: #ffffff;
        margin: 0;
        text-transform: uppercase;
        transition: color 0.3s;
    }
    .header-link:hover .title-text {
        color: #E3000F;
    }
    
    .subtitle-text {
        text-align: center; 
        color: #E3000F; /* SANY Red */
        font-weight: 600;
        letter-spacing: 2px;
        margin-top: 10px;
        font-size: 1.3rem;
    }
    
    /* Global Adjustments */
    .stButton>button {
        font-weight: bold;
        text-transform: uppercase;
        letter-spacing: 1px;
        font-size: 1.1rem;
        padding: 0.75rem 1.5rem;
    }
    h3 {
        text-transform: uppercase;
        letter-spacing: 1px;
        border-bottom: 2px solid #E3000F;
        padding-bottom: 10px;
        margin-top: 20px;
    }

    /* OVERRIDE STREAMLIT ALERTS (Info/Success) TO SANY RED */
    div[data-testid="stAlert"] {
        background-color: rgba(227, 0, 15, 0.1) !important;
        border: 1px solid #E3000F !important;
        color: #ffffff !important;
    }
    div[data-testid="stAlert"] svg {
        fill: #E3000F !important;
    }
</style>
""", unsafe_allow_html=True)

st.markdown("""
<a href="?" target="_self" class="header-link" title="Click to reset App">
    <div class="radar-container">
        <div class="radar"></div>
        <h1 class="title-text">RivalRadar</h1>
    </div>
</a>
""", unsafe_allow_html=True)
st.markdown("<h4 class='subtitle-text'>[ AI-POWERED COMPETITOR ANALYSIS SYSTEM ]</h4>", unsafe_allow_html=True)
st.markdown("---")

MACHINE_TYPES = ["Tracked Excavator", "Wheeled Excavator", "Wheel Loader", "Mining Excavator"]

CATEGORIES = {
    "Tracked Excavator": [
        "SY10U", "SY16C", "SY18U", "SY18C", "SY19E", "SY20C", "SY26U", "SY26C", 
        "SY35U", "SY35C", "SY50U", "SY60U", "SY60C", "SY75C", "SY80U", "SY95C", 
        "SY135C", "SY155U", "SY215C", "SY235C", "SY265C", "SY305C", "SY365H", 
        "SY390H", "SY500H", "SY750H", "SY980H"
    ],
    "Wheeled Excavator": [
        "SY155W", "SY175W", "7t Placeholder", "10t Placeholder"
    ],
    "Wheel Loader": [
        "SW305K", "SW405K"
    ],
    "Mining Excavator": [
        "SY2000H"
    ]
}

PARAMS = {
    "Tracked Excavator": [
        "Engine Make/Model", "Emission Standard", "Gross Power (kW)", "Net Power (kW)", "Max Torque (Nm)", "Displacement (L)", "Number of Cylinders",
        "Operating Weight (kg)", "Transport Weight (kg)", "Counterweight (kg)",
        "Breakout Force - Bucket (kN)", "Tearout Force - Stick (kN)", "Drawbar Pull / Traction Force (kN)", "Max Travel Speed H/L (km/h)", "Swing Speed (rpm)", "Gradeability (%)",
        "Main Pump Type", "Max Flow (l/min)", "System Pressure (bar)", "AUX 1 Flow (l/min)", "AUX 2 Flow (l/min)",
        "Boom Length (mm)", "Standard Stick Length (mm)", "Max Digging Depth (mm)", "Max Digging Reach on Ground (mm)", "Max Dump Height (mm)", "Tail Swing Radius (mm)", "Ground Clearance (mm)", "Track Gauge (mm)", "Track Shoe Width (mm)",
        "Fuel Tank (l)", "Hydraulic System (l)", "Engine Oil (l)", "DEF/AdBlue Tank (l)"
    ],
    "Wheeled Excavator": [
        "Engine Make/Model", "Emission Standard", "Gross Power (kW)", "Net Power (kW)", "Max Torque (Nm)",
        "Operating Weight with Blade (kg)", "Operating Weight with Outriggers (kg)",
        "Breakout Force (kN)", "Tearout Force (kN)", "Swing Speed (rpm)",
        "Transmission Type", "Travel Speed Creep (km/h)", "Travel Speed Low (km/h)", "Travel Speed High (km/h)", "Tire Size/Type", "Steering Radius (mm)", "Front Axle Oscillation Angle (°)",
        "Main Pump Max Flow (l/min)", "System Pressure (bar)", "AUX 1 Flow (l/min)",
        "Max Digging Depth (mm)", "Max Reach (mm)", "Max Dump Height (mm)", "Tail Swing Radius (mm)", "Wheelbase (mm)", "Overall Width (mm)",
        "Fuel Tank (l)", "Hydraulic Tank (l)", "DEF/AdBlue Tank (l)"
    ],
    "Wheel Loader": [
        "Engine Make/Model", "Emission Standard", "Gross Power (kW)", "Max Torque (Nm)", "Displacement (L)",
        "Operating Weight (kg)",
        "Rated Payload (kg)", "Static Tipping Load - Straight (kg)", "Static Tipping Load - Full Turn (kg)", "Bucket Breakout Force (kN)", "Standard Bucket Capacity Heaped (m3)",
        "Raise Time (s)", "Dump Time (s)", "Lower Time (s)", "Total Cycle Time (s)",
        "Transmission Type", "Gears Forward/Reverse", "Max Travel Speed Forward (km/h)", "Articulation Angle (°)", "Turning Radius outside tires (mm)", "Tire Size",
        "Linkage Type (Z-Bar / Parallel)", "Hinge Pin Height at Max Lift (mm)", "Dump Clearance at Max Lift (mm)", "Reach at Max Lift (mm)", "Overall Length (mm)", "Wheelbase (mm)", "Overall Width (mm)",
        "Fuel Tank (l)", "Hydraulic Tank (l)", "DEF/AdBlue Tank (l)"
    ],
    "Mining Excavator": [
        "Engine Make/Model", "Number of Engines", "Gross Power (kW)", "Electric Motor Power (kW)", "Displacement (L)", "Number of Cylinders",
        "Operating Weight - Backhoe (kg)", "Operating Weight - Face Shovel (kg)", "Counterweight (kg)",
        "Standard Bucket Capacity - Backhoe (m3)", "Standard Bucket Capacity - Face Shovel (m3)",
        "Breakout Force (kN)", "Tearout Force (kN)", "Swing Speed (rpm)", "Swing Torque (kNm)", "Traction Force (kN)",
        "Main Pumps Total Flow (l/min)", "System Pressure (bar)", "Swing Circuit Pressure (bar)",
        "Max Digging Depth (mm)", "Max Reach at Ground Level (mm)", "Max Dump Height (mm)", "Max Digging Height (mm)",
        "Track Gauge (mm)", "Track Shoe Width (mm)", "Overall Width (mm)", "Tail Swing Radius (mm)", "Ground Clearance (mm)", "Number of Track Rollers",
        "Fuel Tank (l)", "Hydraulic System (l)", "Cooling System (l)", "Engine Oil (l)"
    ]
}

DB_FILE = "machine_database.json"

def load_db():
    if os.path.exists(DB_FILE):
        try:
            with open(DB_FILE, "r") as f:
                return json.load(f)
        except:
            return {}
    return {}

def save_db(data):
    with open(DB_FILE, "w") as f:
        json.dump(data, f, indent=4)
        
    if "GITHUB_TOKEN" in st.secrets and "GITHUB_REPO" in st.secrets:
        try:
            from github import Github
            g = Github(st.secrets["GITHUB_TOKEN"])
            repo = g.get_repo(st.secrets["GITHUB_REPO"])
            json_str = json.dumps(data, indent=4)
            try:
                contents = repo.get_contents(DB_FILE)
                repo.update_file(contents.path, "Auto-Sync: Database updated", json_str, contents.sha)
            except:
                repo.create_file(DB_FILE, "Auto-Sync: Database created", json_str)
        except Exception as e:
            print(f"GitHub Sync Error: {e}")

db = load_db()

def extract_number(val):
    if isinstance(val, (int, float)):
        return float(val)
    if isinstance(val, str):
        match = re.search(r"[-+]?\d*\.\d+|\d+", val.replace(',', '.'))
        if match:
            return float(match.group())
    return None

def get_template_path():
    if os.path.exists("sany_template.pptx"):
        return "sany_template.pptx"
    for root, dirs, files in os.walk("."):
        for file in files:
            if "sany" in file.lower() and file.endswith(".pptx") and not file.startswith("._") and not file.startswith("."):
                return os.path.join(root, file)
    return None

def custom_youtube_search(query, time_filter="Any Time", limit=5):
    # Determine the time filter parameter for YouTube
    sp_param = ""
    if time_filter == "Last 12 Months":
        sp_param = "&sp=EgQIBBAB"  
    elif time_filter == "Last 30 Days":
        sp_param = "&sp=EgQIBRAB"  

    url = f"https://www.youtube.com/results?search_query={urllib.parse.quote(query)}{sp_param}"
    
    try:
        req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'})
        html = urllib.request.urlopen(req).read().decode('utf-8')
        
        # Extract Video IDs using Regex
        video_ids = re.findall(r"\"videoId\":\"([a-zA-Z0-9_-]{11})\"", html)
        unique_ids = list(dict.fromkeys(video_ids)) 
        
        results = []
        for vid in unique_ids[:limit]:
            title_match = re.search(r'"videoId":"' + vid + r'".*?"title":\{"runs":\[\{"text":"([^"]+)"', html)
            title = title_match.group(1) if title_match else f"YouTube Video ({vid})"
            
            results.append({
                'id': vid,
                'title': title,
                'link': f"https://www.youtube.com/watch?v={vid}"
            })
        return results
    except Exception as e:
        print(f"Custom YouTube search failed: {e}")
        return []

def create_pitch_deck(baseline_name, competitor_names, ai_text, df_battle):
    template_path = get_template_path()
    if not template_path:
        prs = Presentation() 
    else:
        try: prs = Presentation(template_path)
        except: prs = Presentation()
        
    # Title Slide
    try:
        if len(prs.slides) > 0:
            slide1 = prs.slides[0]
            replaced_title = False
            for shape in slide1.shapes:
                if shape.has_text_frame and "XXX" in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace("XXX", f"Tactical Product Comparison\n{baseline_name} vs. {', '.join(competitor_names)}")
                    replaced_title = True
                    break
            if not replaced_title:
                text_shapes_1 = [s for s in slide1.shapes if s.has_text_frame]
                text_shapes_1.sort(key=lambda s: s.width * s.height if s.width and s.height else 0, reverse=True)
                if len(text_shapes_1) > 0:
                    text_shapes_1[0].text_frame.text = f"Tactical Product Comparison\n{baseline_name} vs. {', '.join(competitor_names)}"
    except: pass
        
    # Table Slide
    try:
        table_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        table_slide = prs.slides.add_slide(table_layout)
        slide_id_list = prs.slides._sldIdLst
        new_slide_id = slide_id_list[-1]
        slide_id_list.remove(new_slide_id)
        slide_id_list.insert(1, new_slide_id)
        
        for shape in table_slide.shapes:
            if shape.has_text_frame and "XXXX" in shape.text_frame.text:
                shape.text_frame.text = shape.text_frame.text.replace("XXXX", "Technical Specifications Matrix")
            elif shape.has_text_frame and "XXX" in shape.text_frame.text:
                shape.text_frame.text = ""

        if table_slide.shapes.title and "Technical" not in table_slide.shapes.title.text:
            table_slide.shapes.title.text = "Technical Specifications Matrix"

        rows, cols = len(df_battle.columns), len(df_battle) + 1
        left, top, width, height = Inches(0.5), Inches(1.8), Inches(9.0), Inches(0.4 * rows) 
        shape = table_slide.shapes.add_table(rows, cols, left, top, width, height)
        table = shape.table
        
        table.columns[0].width = Inches(2.5)
        for i in range(1, cols): table.columns[i].width = Inches(6.5 / (cols - 1))
            
        for row_idx, col_name in enumerate(df_battle.columns):
            cell = table.cell(row_idx, 0)
            cell.text = str(col_name)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            
        for m_idx, row_data in df_battle.iterrows():
            col_idx = m_idx + 1
            for r_idx, col_name in enumerate(df_battle.columns):
                cell = table.cell(r_idx, col_idx)
                val = str(row_data[col_name])
                if val == "nan" or val == "None": val = "-"
                cell.text = val
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                if r_idx == 0: cell.text_frame.paragraphs[0].font.bold = True
    except: pass

    # AI Analysis Slide
    if ai_text:
        try:
            if len(prs.slides) > 2: slide3 = prs.slides[2]
            else: slide3 = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
            
            tf_content = None
            for shape in slide3.shapes:
                if shape.has_text_frame:
                    if "XXXX" in shape.text_frame.text: shape.text_frame.text = shape.text_frame.text.replace("XXXX", "Competitive Analysis & Pitch")
                    elif "XXX" in shape.text_frame.text: tf_content = shape.text_frame
            
            if not tf_content:
                text_shapes_3 = [s for s in slide3.shapes if s.has_text_frame]
                text_shapes_3.sort(key=lambda s: s.width * s.height if s.width and s.height else 0, reverse=True)
                if len(text_shapes_3) > 0: tf_content = text_shapes_3[0].text_frame
                
            if slide3.shapes.title and "Competitive" not in slide3.shapes.title.text: slide3.shapes.title.text = "Competitive Analysis & Pitch"
                
            if tf_content:
                tf_content.text = "" 
                tf_content.word_wrap = True
                try: tf_content.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                except: pass
                
                lines = ai_text.split('\n')
                first_paragraph = True
                for line in lines:
                    line = line.strip()
                    if not line or line.startswith('|') or line.startswith('---') or line.startswith('='): continue
                    clean_line = line.replace('**', '').replace('### ', '').replace('## ', '')
                    p = tf_content.paragraphs[0] if first_paragraph else tf_content.add_paragraph()
                    first_paragraph = False
                    if clean_line.startswith('* ') or clean_line.startswith('- '):
                        p.text = clean_line[2:]
                        p.level, p.font.size = 1, Pt(14)
                    else:
                        p.text = clean_line
                        p.level, p.font.bold, p.font.size = 0, True, Pt(16)
        except: pass
            
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

def create_video_intel_deck(target_machine, videos_data, exec_summary):
    template_path = get_template_path()
    if not template_path: prs = Presentation() 
    else:
        try: prs = Presentation(template_path)
        except: prs = Presentation()

    content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    # Clean existing dummy slides if using template
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    if slide.shapes.title: slide.shapes.title.text = "Video Intelligence Report"
    if len(slide.placeholders) > 1: slide.placeholders[1].text = f"Target Machine: {target_machine.upper()}\nGlobal Operator Sentiment"

    # Slides 2-6: Individual Videos
    for vid in videos_data:
        slide = prs.slides.add_slide(content_layout)
        if slide.shapes.title: slide.shapes.title.text = vid['title'][:50] + "..."
        
        tf = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                tf = shape.text_frame
                break
        if not tf:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
            tf = txBox.text_frame

        tf.text = f"Video Link: {vid['link']}\n\n"
        tf.word_wrap = True
        try: tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except: pass

        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        
        # Add AI Summary Points (Cleaned from Markdown)
        lines = vid['summary'].split('\n')
        for line in lines:
            line = line.strip().replace('**', '').replace('### ', '').replace('## ', '').replace('# ', '')
            if line:
                p = tf.add_paragraph()
                if line.startswith('* ') or line.startswith('- '):
                    p.text = line[2:] 
                    p.level = 1
                else:
                    p.text = line
                    p.level = 0
                p.font.size = Pt(16)

    # Final Slide: Executive Summary
    slide = prs.slides.add_slide(content_layout)
    if slide.shapes.title: slide.shapes.title.text = "Executive Summary & Strategy"
    
    tf = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            tf = shape.text_frame
            break
    if not tf:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        tf = txBox.text_frame

    tf.text = ""
    tf.word_wrap = True
    try: tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except: pass
    
    lines = exec_summary.split('\n')
    for line in lines:
        line = line.strip().replace('**', '').replace('### ', '').replace('## ', '').replace('# ', '')
        if line:
            p = tf.add_paragraph()
            if line.startswith('* ') or line.startswith('- '):
                p.text = line[2:] 
                p.level = 1
            else:
                p.text = line
                p.level = 0
                
            p.font.size = Pt(16)
            if "The Good" in line or "The Bad" in line or "SANY" in line:
                p.font.bold = True

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

if "custom_params" not in st.session_state:
    st.session_state.custom_params = {}
    for mt in MACHINE_TYPES:
        st.session_state.custom_params[mt] = []

if "ai_analysis_cache" not in st.session_state: st.session_state.ai_analysis_cache = ""
if "video_intel_videos" not in st.session_state: st.session_state.video_intel_videos = []
if "video_intel_summary" not in st.session_state: st.session_state.video_intel_summary = ""
if "video_intel_machine" not in st.session_state: st.session_state.video_intel_machine = ""

st.sidebar.markdown("### NAVIGATION")
app_mode = st.sidebar.radio("Navigate to:", ["Scanner", "Database", "Product Comparison", "News Radar", "Video Intelligence"])
st.sidebar.markdown("---")

if "GEMINI_API_KEY" in st.secrets and st.secrets["GEMINI_API_KEY"]:
    api_key = st.secrets["GEMINI_API_KEY"]
    st.sidebar.error("API-Key loaded successfully.") 
else:
    st.sidebar.markdown("### SYSTEM SETTINGS")
    api_key = st.sidebar.text_input("Gemini API Key (Manual)", type="password")

selected_model_name = "gemini-3.1-flash-lite"
if api_key:
    genai.configure(api_key=api_key)
    st.sidebar.markdown("### AI MODEL")
    try:
        available_models = [m.name.replace("models/", "") for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
        if available_models:
            default_idx = 0
            if "gemini-3.1-flash-lite" in available_models: default_idx = available_models.index("gemini-3.1-flash-lite")
            elif any("flash-lite" in m for m in available_models): default_idx = next(i for i, m in enumerate(available_models) if "flash-lite" in m)
            selected_model_name = st.sidebar.selectbox("Active Model:", available_models, index=default_idx)
    except:
        selected_model_name = st.sidebar.text_input("Manual Model Input:", value="gemini-3.1-flash-lite")

st.sidebar.markdown("---")

if app_mode == "Scanner":
    st.markdown("### UPLOAD DATASHEETS")
    selected_machine_type = st.radio("Configure AI Brain For:", MACHINE_TYPES, horizontal=True)
    
    current_categories = CATEGORIES[selected_machine_type]
    current_base_params = PARAMS[selected_machine_type]
    all_available_params = current_base_params + st.session_state.custom_params[selected_machine_type]

    with st.expander(f"CONFIGURE PARAMS ({selected_machine_type})", expanded=False):
        st.error(f"The system recognizes {len(all_available_params)} parameters for {selected_machine_type}.")
        new_param_input = st.text_input("Add custom parameter:", placeholder="e.g., Track width (mm)", key=f"new_param_{selected_machine_type}")
        if st.button("Save Parameter", use_container_width=True):
            if new_param_input and new_param_input.strip() not in all_available_params:
                st.session_state.custom_params[selected_machine_type].append(new_param_input.strip())
                st.rerun()
    
    uploaded_files = st.file_uploader("Drop brochures or datasheets here (PDF, PNG, JPG)", type=["pdf", "png", "jpg", "jpeg"], accept_multiple_files=True)

    machine_configs = {}
    if uploaded_files:
        if not api_key: st.error("ACCESS DENIED: Please provide a valid API Key.")
        else:
            for file in uploaded_files:
                sub_col1, sub_col2 = st.columns(2)
                with sub_col1:
                    m_name = st.text_input(f"Machine Name:", value=file.name.rsplit('.', 1)[0], key=f"name_{file.name}")
                with sub_col2:
                    m_cat = st.selectbox(f"Baseline Class ({selected_machine_type}):", options=current_categories, key=f"cat_{file.name}")
                machine_configs[file.name] = {"name": m_name, "category": m_cat}

    if st.button("INITIATE AI SCAN (EXTRACT ALL DATA)", type="primary", use_container_width=True):
        if uploaded_files:
            progress_bar = st.progress(0)
            status_text = st.empty()
            for index, file in enumerate(uploaded_files):
                current_machine = machine_configs[file.name]["name"]
                current_category = machine_configs[file.name]["category"]
                status_text.text(f"Scanning '{current_machine}'...")
                
                file.seek(0)
                file_part = {"mime_type": file.type, "data": file.read()}
                prompt = f"""You are a precise technical data extraction assistant analyzing a {selected_machine_type}. Focus EXCLUSIVELY on: "{current_machine}". Extract exact values for: {json.dumps(all_available_params)}. 1. Valid JSON object only. 2. Exact keys. 3. Use "?" if missing. 4. Convert imperial to metric. ONLY JSON."""
                
                try:
                    model = genai.GenerativeModel(selected_model_name)
                    response = model.generate_content(contents=[file_part, prompt], generation_config={"response_mime_type": "application/json"})
                    raw_text = response.text.strip()
                    marker = "`" * 3
                    if raw_text.startswith(marker + "json"): raw_text = raw_text[7:]
                    elif raw_text.startswith(marker): raw_text = raw_text[3:]
                    if raw_text.endswith(marker): raw_text = raw_text[:-3]
                        
                    extracted_json = json.loads(raw_text.strip())
                    extracted_json["Machine"], extracted_json["Category"], extracted_json["Machine Type"] = current_machine, current_category, selected_machine_type
                    db[f"{current_machine} ({current_category})"] = extracted_json
                    save_db(db)
                except Exception as e:
                    st.error(f"Scan failed for '{current_machine}'. Error: {e}")
                
                progress_bar.progress((index + 1) / len(uploaded_files))
            
            status_text.text("Data extraction complete. Machines synced to database.")
            time.sleep(2)
            st.rerun()

elif app_mode == "Database":
    st.markdown("### MACHINE DATABASE (LIVE EDITOR)")
    if db:
        col_filter_type, col_filter_cat, col_delete = st.columns(3)
        all_types = list(set([v.get("Machine Type", "Tracked Excavator") for v in db.values()]))
        with col_filter_type: type_filter = st.selectbox("Filter by Machine Type:", ["All"] + all_types)
        with col_filter_cat:
            cat_options = ["All"]
            if type_filter != "All": cat_options += CATEGORIES.get(type_filter, [])
            else:
                for cats in CATEGORIES.values(): cat_options += cats
            cat_filter = st.selectbox("Filter by SANY Class:", cat_options)
        with col_delete:
            to_delete = st.selectbox("Remove faulty record:", options=["None"] + list(db.keys()))
            if st.button("DELETE RECORD", use_container_width=True):
                if to_delete != "None":
                    del db[to_delete]
                    save_db(db)
                    st.rerun()
                    
        st.markdown("---")
        df_lib = pd.DataFrame(list(db.values()))
        if 'Category' in df_lib.columns: df_lib = df_lib.rename(columns={'Category': 'Class'})
        if 'Machine Type' not in df_lib.columns: df_lib['Machine Type'] = "Tracked Excavator"
            
        cols = ['Machine', 'Machine Type', 'Class'] + [c for c in df_lib.columns if c not in ['Machine', 'Machine Type', 'Class']]
        df_lib = df_lib[cols]
        if type_filter != "All": df_lib = df_lib[df_lib['Machine Type'] == type_filter]
        if cat_filter != "All": df_lib = df_lib[df_lib['Class'] == cat_filter]
            
        edited_df = st.data_editor(df_lib, use_container_width=True, num_rows="dynamic")
        if st.button("SAVE DATABASE CHANGES", type="primary", use_container_width=True):
            new_db = {}
            for _, row in edited_df.iterrows():
                new_key = f"{row['Machine']} ({row['Class'] if 'Class' in row else row.get('Category', '')})"
                row_dict = row.to_dict()
                if 'Class' in row_dict: row_dict['Category'] = row_dict.pop('Class')
                new_db[new_key] = {k: v for k, v in row_dict.items() if pd.notna(v)}
            save_db(new_db)
            db.clear(); db.update(new_db)
            st.error("Database successfully updated and synced!")
            time.sleep(1)
            st.rerun()
        
        st.markdown("---")
        excel_buffer_lib = io.BytesIO()
        with pd.ExcelWriter(excel_buffer_lib, engine='openpyxl') as writer: df_lib.to_excel(writer, index=False)
        st.download_button("DOWNLOAD DATABASE (.xlsx)", excel_buffer_lib.getvalue(), "Database.xlsx", use_container_width=True)
    else: st.error("The database is currently empty.")

elif app_mode == "Product Comparison":
    st.markdown("### PRODUCT COMPARISON (THE ARENA)")
    if not db: st.error("No data available yet.")
    else:
        arena_type = st.radio("Filter Arena By Machine Type:", MACHINE_TYPES, horizontal=True)
        filtered_db_keys = [k for k, v in db.items() if v.get("Machine Type", "Tracked Excavator") == arena_type]
        if not filtered_db_keys: st.error(f"No machines found for '{arena_type}' in the database.")
        else:
            arena_col1, arena_col2 = st.columns(2)
            with arena_col1: baseline_sel = st.selectbox("Own Product (SANY Baseline):", options=filtered_db_keys)
            with arena_col2: competitors_sel = st.multiselect("Competitor Models:", options=[k for k in filtered_db_keys if k != baseline_sel])
                
            st.markdown("---")
            all_available_params = PARAMS[arena_type] + st.session_state.custom_params[arena_type]
            
            if arena_type == "Tracked Excavator": default_params = ["Operating Weight (kg)", "Net Power (kW)", "Max Digging Depth (mm)", "Breakout Force - Bucket (kN)", "AUX 1 Flow (l/min)"]
            elif arena_type == "Wheeled Excavator": default_params = ["Operating Weight with Blade (kg)", "Net Power (kW)", "Max Travel Speed High (km/h)", "Breakout Force (kN)", "Tail Swing Radius (mm)"]
            elif arena_type == "Mining Excavator": default_params = ["Operating Weight - Backhoe (kg)", "Gross Power (kW)", "Standard Bucket Capacity - Backhoe (m3)", "Breakout Force (kN)", "Max Digging Depth (mm)"]
            else: default_params = ["Operating Weight (kg)", "Rated Payload (kg)", "Static Tipping Load - Full Turn (kg)", "Standard Bucket Capacity Heaped (m3)", "Total Cycle Time (s)"]
                
            compare_params = st.multiselect("Select parameters for Matrix & AI Analysis:", options=all_available_params, default=[p for p in default_params if p in all_available_params])
            ai_persona = st.radio("Select AI Persona:", ["Sales Mode (Punchy, Strategic, ROI & Sales Focus)", "R&D Mode (Technical, Analytical, Engineering & Structure)"])
            pwr_charts = st.checkbox("Visual Performance Comparison (Charts in App)")
            pwr_ampel = st.checkbox("Strengths/Weaknesses Profile (For Pitch Deck)")
            pwr_pitch = st.checkbox("Sales/Tech Arguments (For Pitch Deck)")
            
            if st.button("GENERATE COMPARISON", type="primary", use_container_width=True):
                if competitors_sel and compare_params:
                    battle_data = [db[baseline_sel]] + [db[k] for k in competitors_sel]
                    df_battle_full = pd.DataFrame(battle_data).drop(columns=['Category', 'Machine Type', 'Class'], errors='ignore')
                    df_battle_filtered = df_battle_full[['Machine'] + [p for p in compare_params if p in df_battle_full.columns]]
                    st.session_state.current_df_battle = df_battle_filtered
                    
                    st.markdown("---")
                    st.write("### Raw Data Matrix (Filtered)")
                    st.dataframe(df_battle_filtered.set_index("Machine").T, use_container_width=True)
                    
                    if pwr_charts:
                        st.markdown("---")
                        st.write("### Visual Performance Comparison")
                        chart_cols = st.columns(2)
                        col_idx = 0
                        for metric in compare_params:
                            if metric in df_battle_filtered.columns:
                                chart_data = [{"Machine": row['Machine'], metric: extract_number(row.get(metric))} for _, row in df_battle_filtered.iterrows() if extract_number(row.get(metric)) is not None]
                                if chart_data:
                                    chart = alt.Chart(pd.DataFrame(chart_data)).mark_bar().encode(
                                        x=alt.X('Machine', title=None, axis=alt.Axis(labelAngle=-45)),
                                        y=alt.Y(metric, title=None), color=alt.Color('Machine', legend=None), tooltip=['Machine', metric]
                                    ).properties(height=300, title=metric)
                                    with chart_cols[col_idx % 2]: st.altair_chart(chart, use_container_width=True)
                                    col_idx += 1

                    if pwr_ampel or pwr_pitch:
                        st.markdown("---")
                        st.write("### AI Competitive Analysis")
                        with st.spinner("The AI is analyzing the data..."):
                            baseline_name = db[baseline_sel]['Machine']
                            competitor_names = [db[k]['Machine'] for k in competitors_sel]
                            sys_prompt = f"You are a {'Senior Sales Strategist' if 'Sales' in ai_persona else 'Senior R&D Engineer'} evaluating {arena_type} models. English only. Baseline: '{baseline_name}'. Competitors: {', '.join(competitor_names)}.\nData: {df_battle_filtered.to_dict(orient='records')}\nCRITICAL: NEVER USE MARKDOWN TABLES! Write plain text paragraphs and short bullet points starting with '*'.\n"
                            if pwr_ampel: sys_prompt += "- Objective assessment (Superior, Tie, Competitor superior) formatted as SHORT text bullets.\n"
                            if pwr_pitch: sys_prompt += "- Short, punchy arguments (Max 1 sentence each).\n"
                            try:
                                model = genai.GenerativeModel(selected_model_name)
                                response = model.generate_content(sys_prompt)
                                st.session_state.ai_analysis_cache = response.text
                                st.markdown(response.text)
                            except Exception as e: st.error(f"AI Analysis failed: {e}")

            if st.session_state.get("ai_analysis_cache") and competitors_sel and "current_df_battle" in st.session_state:
                st.markdown("---")
                baseline_name = db[baseline_sel]['Machine']
                ppt_file = create_pitch_deck(baseline_name, [db[k]['Machine'] for k in competitors_sel], st.session_state.ai_analysis_cache, st.session_state.current_df_battle)
                st.download_button("DOWNLOAD PITCH DECK (.pptx)", data=ppt_file.getvalue(), file_name=f"SANY_Pitch_{baseline_name}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True, type="primary")

elif app_mode == "News Radar":
    st.markdown("### 📡 OEM NEWS RADAR")
    oem_list = ["Caterpillar", "Komatsu", "Volvo Construction Equipment", "Liebherr", "Develon", "Yanmar", "Kubota", "Hitachi Construction Machinery"]
    col1, col2 = st.columns([2, 1])
    with col1: selected_oem = st.selectbox("Target OEM:", ["All Competitors"] + oem_list)
    with col2: search_focus = st.selectbox("Focus Area:", ["Excavator", "Wheel Loader", "Corporate / Tech"])
        
    button_text = "INITIATE RADAR SWEEP FOR ALL" if selected_oem == "All Competitors" else f"INITIATE RADAR SWEEP FOR {selected_oem.upper()}"
    if st.button(button_text, type="primary", use_container_width=True):
        if not api_key: st.error("ACCESS DENIED: Please provide a valid API Key to use the AI Briefing.")
        else:
            with st.spinner("Intercepting global feeds and generating AI summary..."):
                oem_query_part = "(" + " OR ".join([f'"{o}"' for o in oem_list]) + ")" if selected_oem == "All Competitors" else f'"{selected_oem}"'
                query = f'{oem_query_part} (equipment OR machinery) (launch OR technology OR "press release")' if search_focus == "Corporate / Tech" else f'{oem_query_part} "{search_focus}" (launch OR reveal OR "press release" OR new)'
                try:
                    feed = feedparser.parse(f"https://news.google.com/rss/search?q={urllib.parse.quote(query)}&hl=en-US&gl=US&ceid=US:en")
                    if not feed.entries: st.error(f"No recent high-impact news found.")
                    else:
                        news_text = ""
                        st.markdown("---")
                        result_limit = 10 if selected_oem == "All Competitors" else 5
                        st.markdown(f"#### INTERCEPTED SIGNALS (TOP {result_limit})")
                        for entry in feed.entries[:result_limit]:
                            st.markdown(f"**[{entry.title}]({entry.link})**")
                            st.caption(f"Broadcast Date: {entry.published if hasattr(entry, 'published') else 'Recent'}")
                            news_text += f"- {entry.title}\n"
                        
                        st.markdown("---")
                        st.markdown("#### 🧠 AI STRATEGY BRIEFING")
                        target_text = "our top competitors in the global market" if selected_oem == "All Competitors" else f"our competitor {selected_oem}"
                        sys_prompt = f"You are a Senior Market Intelligence Analyst for SANY Europe. Analyze these recent headlines about {target_text}:\n{news_text}\nProvide a punchy, professional 3-sentence executive summary in English. Highlight any potential strategic threats, market trends, or product advancements that SANY sales teams need to watch out for. No markdown tables, just plain text."
                        try:
                            model = genai.GenerativeModel(selected_model_name)
                            st.markdown(model.generate_content(sys_prompt).text)
                        except Exception as e: st.error(f"AI Briefing failed: {e}")
                except Exception as e: st.error(f"Radar sweep failed. Signal lost: {e}")

elif app_mode == "Video Intelligence":
    st.markdown("### 🎬 VIDEO INTELLIGENCE (OPERATOR SENTIMENT)")
    st.error("Extract unfiltered operator opinions from YouTube. The AI pulls transcripts from the top 5 videos and auto-generates a presentation slide deck.")

    col1, col2, col3 = st.columns([2, 1, 1])
    with col1:
        target_machine = st.text_input("Target Machine / Competitor (e.g., CAT 320, Liebherr R 920):", placeholder="Enter machine model...")
    with col2:
        region_filter = st.selectbox("Region / Language:", [
            "World Wide (English)", "European Market (Mixed)", "Germany (German)", "France (French)", "Italy (Italian)"
        ])
    with col3:
        time_filter = st.selectbox("Upload Date:", [
            "Any Time", "Last 12 Months", "Last 30 Days"
        ])

    if st.button("GENERATE VIDEO PPT DECK", type="primary", use_container_width=True):
        if not api_key:
            st.error("ACCESS DENIED: Please provide a valid API Key.")
        elif not target_machine:
            st.warning("Please enter a target machine to scan.")
        else:
            status_placeholder = st.empty()
            progress_bar = st.progress(0)
            
            # Setup search terms based on region
            lang_query = "review"
            if "German" in region_filter: lang_query = "testbericht OR erfahrungen"
            elif "French" in region_filter: lang_query = "avis OR essai"
            elif "Italian" in region_filter: lang_query = "recensione OR prova"
            elif "European" in region_filter: lang_query = "walkaround"

            try:
                search_term = f"{target_machine} {lang_query} excavator"
                status_placeholder.info(f"Phase 1: Direct YouTube bypass to find top 5 videos on '{target_machine}' ({time_filter})...")
                
                results = custom_youtube_search(search_term, time_filter, limit=5)
                
                if not results:
                    status_placeholder.error("No videos found for this machine. Try a broader search term.")
                else:
                    model = genai.GenerativeModel(selected_model_name)
                    videos_data = []
                    full_transcripts_for_exec = ""
                    
                    for idx, vid in enumerate(results):
                        status_placeholder.info(f"Phase 2: Extracting and summarizing transcript for Video {idx+1}/5...")
                        vid_id = vid['id']
                        title = vid['title']
                        link = vid['link']
                        
                        try:
                            # Safest method across all library versions - brute forcing the most common languages
                            transcript_data = YouTubeTranscriptApi.get_transcript(
                                vid_id, 
                                languages=['en', 'en-US', 'en-GB', 'de', 'fr', 'es', 'it', 'nl', 'ru', 'zh', 'ja', 'ko', 'pt', 'tr', 'pl', 'sv', 'ar']
                            )
                                
                            transcript_text = " ".join([t['text'] for t in transcript_data])
                            
                            # AI summarize this single video
                            prompt = f"Summarize this YouTube video transcript about the construction machine '{target_machine}'. Extract exactly 3 short, punchy bullet points highlighting the operator's opinion (pros/cons). English only. Start each point with '* '. NO MARKDOWN HEADERS. Transcript: {transcript_text[:8000]}"
                            vid_summary = model.generate_content(prompt).text
                            full_transcripts_for_exec += f"\n\nVideo: {title}\nSummary: {vid_summary}"
                            
                        except Exception as e:
                            # Bulletproof Fallback if no captions exist
                            vid_summary = f"* AI Note: Transcript extraction failed.\n* The video might have no spoken words or captions are disabled.\n* Recommend manual review by sales rep."
                            
                        videos_data.append({'title': title, 'link': link, 'summary': vid_summary})
                        progress_bar.progress((idx + 1) / 6) # 6 steps total

                    # --- Generate Executive Summary ---
                    status_placeholder.info("Phase 3: Generating Global Executive Summary...")
                    exec_prompt = f"""
                    You are a Market Intelligence Analyst for SANY Europe. 
                    Based on these 5 video summaries regarding the {target_machine}, create an executive summary.
                    
                    {full_transcripts_for_exec}
                    
                    Format exactly with these headers (English only, bullet points):
                    * Operator Praises (The Good):
                    * Operator Complaints (The Bad):
                    * SANY Sales Counter-Argument: (1 punchy sentence how a SANY rep can use these weaknesses to sell a SANY)
                    
                    CRITICAL: Do not use any markdown formatting like ### or **. Just pure text and bullet points.
                    """
                    exec_summary_text = model.generate_content(exec_prompt).text
                    progress_bar.progress(1.0)
                    
                    # Save to session state
                    st.session_state.video_intel_videos = videos_data
                    st.session_state.video_intel_summary = exec_summary_text
                    st.session_state.video_intel_machine = target_machine
                    
                    status_placeholder.success("✅ Analysis Complete! Pitch Deck is ready.")

            except Exception as e:
                st.error(f"Process failed: {e}")

    # Display Report & Download Button if cache exists
    if st.session_state.get("video_intel_summary"):
        st.markdown("---")
        st.markdown("#### 🧠 AI EXECUTIVE SUMMARY (PREVIEW)")
        st.markdown(st.session_state.video_intel_summary)
        
        st.markdown("---")
        st.markdown("#### 📥 DOWNLOAD REPORT")
        
        ppt_file = create_video_intel_deck(
            st.session_state.video_intel_machine, 
            st.session_state.video_intel_videos, 
            st.session_state.video_intel_summary
        )
        
        st.download_button(
            label=f"DOWNLOAD VIDEO INTELLIGENCE DECK (.pptx)",
            data=ppt_file.getvalue(),
            file_name=f"SANY_Video_Intel_{st.session_state.video_intel_machine.replace(' ', '_')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
            type="primary"
        )
